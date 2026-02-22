"""
DocuForge AI — Multimodal Parser

Extracts content from various document formats including PDFs, images
with OCR, and audio files with transcription. Returns normalized text
suitable for downstream processing.
"""

import base64
import logging
import time
from pathlib import Path

import pandas as pd
import yaml
from langchain_core.messages import HumanMessage
from pptx import Presentation

import whisper

from core.llm_router import get_llm

logger = logging.getLogger(__name__)

# Cache Whisper model at module level to avoid reloading
_whisper_model = None


def _load_config() -> dict:
    """Load configuration from docuforge_config.yaml."""
    config_path = Path(__file__).parent.parent.parent / "config" / "docuforge_config.yaml"
    with open(config_path) as f:
        return yaml.safe_load(f)


def _get_whisper_model():
    """Get or load the Whisper model, cached at module level."""
    global _whisper_model
    if _whisper_model is None:
        logger.info("Loading Whisper base model (one-time cache)")
        _whisper_model = whisper.load_model("base")
    return _whisper_model


def parse_pdf(file_path: str) -> str:
    """
    Extract all text from a PDF file using PyMuPDF (fitz).
    Returns concatenated text from all pages. Handles encrypted PDFs gracefully.
    
    Args:
        file_path: Path to PDF file
        
    Returns:
        Extracted text or error message
    """
    try:
        import fitz
    except ImportError:
        logger.error("PyMuPDF (fitz) not installed")
        return "PDF could not be processed (fitz not available)."

    try:
        text_parts = []
        page_count = 0
        with fitz.open(file_path) as pdf:
            # Check if encrypted
            if pdf.is_encrypted:
                logger.warning("PDF is encrypted, cannot extract text")
                return "PDF is encrypted and cannot be read without a password."
            
            page_count = len(pdf)
            for page_num, page in enumerate(pdf):
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"\n--- Page {page_num + 1} ---\n{text}")

        if not text_parts:
            logger.warning("No text extracted from PDF")
            return "No text content found in PDF."

        result = "\n".join(text_parts)
        logger.info("Parsed PDF: %d pages, %d chars extracted", page_count, len(result))
        return result

    except Exception as e:
        logger.error("PDF parsing failed: %s", e)
        return f"PDF could not be processed: {str(e)}"


def parse_image(file_path: str) -> str:
    """
    Extract text and visual content from an image using the configured LLM
    with multimodal (vision) capability.
    
    Args:
        file_path: Path to image file
        
    Returns:
        Extracted text or error message
    """
    try:
        with open(file_path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode("utf-8")

        file_ext = Path(file_path).suffix.lower()
        mime_type_map = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }
        mime_type = mime_type_map.get(file_ext, "image/jpeg")

        llm = get_llm(task_type="multimodal", has_image=True)

        message = HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime_type};base64,{image_data}"},
                },
                {
                    "type": "text",
                    "text": "Extract and describe all text, data, and visual content from this image in detail.",
                },
            ]
        )

        result = llm.invoke([message])
        text = result.content if isinstance(result.content, str) else str(result.content)
        logger.info("Parsed image: extracted %d chars", len(text))
        return text

    except Exception as e:
        logger.error("Image parsing failed: %s", e)
        return "Image could not be processed."


def parse_audio(file_path: str) -> str:
    """
    Transcribe audio file to text using Whisper.
    Caches model at module level for performance.
    
    Args:
        file_path: Path to audio file
        
    Returns:
        Transcribed text or error message
    """
    try:
        start_time = time.time()
        model = _get_whisper_model()
        result = model.transcribe(file_path)
        transcript = result.get("text", "").strip()
        
        if not transcript:
            logger.warning("Audio transcription returned empty result")
            return "Audio transcription failed (no text detected)."
        
        elapsed = time.time() - start_time
        logger.info("Transcribed audio in %.1f seconds: %d chars", elapsed, len(transcript))
        return transcript

    except Exception as e:
        logger.error("Audio transcription failed: %s", e)
        return "Audio transcription failed."


def parse_excel(file_path: str) -> str:
    """
    Extract data from Excel file using pandas.
    Reads all sheets and converts each to string format.
    
    Args:
        file_path: Path to Excel file
        
    Returns:
        Extracted text or error message
    """
    try:
        # Try Excel first, fallback to CSV
        try:
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
        except Exception:
            logger.info("Not an Excel file, attempting CSV parse")
            df = pd.read_csv(file_path)
            if df.empty:
                return "No data found in spreadsheet."
            return df.to_string(index=False)

        sheet_texts = []
        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if df.empty:
                logger.debug("Skipping empty sheet: %s", sheet_name)
                continue
            
            sheet_str = df.to_string(index=False)
            sheet_texts.append(f"=== Sheet: {sheet_name} ===\n{sheet_str}")

        if not sheet_texts:
            logger.warning("No data found in any Excel sheet")
            return "No data found in spreadsheet."

        result = "\n\n".join(sheet_texts)
        logger.info("Parsed Excel: %d non-empty sheets, %d chars", len(sheet_texts), len(result))
        return result

    except Exception as e:
        logger.error("Excel parsing failed: %s", e)
        return "Excel file could not be processed."


def parse_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint presentation.
    Returns all text found in text frames across all slides.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        Extracted text or error message
    """
    try:
        prs = Presentation(file_path)
        slide_texts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_content.append(text)

            if slide_content:
                slide_text = "\n".join(slide_content)
                slide_texts.append(f"=== Slide {slide_num + 1} ===\n{slide_text}")

        if not slide_texts:
            logger.warning("No text content found in presentation")
            return "No text content found in presentation."

        result = "\n\n".join(slide_texts)
        logger.info("Parsed PPTX: %d slides with text, %d chars", len(slide_texts), len(result))
        return result

    except Exception as e:
        logger.error("PPTX parsing failed: %s", e)
        return "PowerPoint file could not be processed."

